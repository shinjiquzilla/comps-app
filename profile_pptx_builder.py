"""
Profile PPTX Builder — python-pptx でCompany Profile PowerPointを生成。

テンプレート仕様（牧野フライス v3準拠）:
  - スライドサイズ: 10" x 7.5"（標準ワイドスクリーン）
  - フォント: Meiryo（日本語）/ Arial（英語）
  - カラー: #4472C4（セクションバー青）、#333333（本文）
  - 生成スライド: Overview, Directors, Comps, Financial Analysis
"""

import io
import os
import re
import tempfile
from datetime import date
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.util import Inches, Pt, Emu

# ---------------------------------------------------------------------------
# Constants — テンプレート準拠のスタイル定数
# ---------------------------------------------------------------------------

# スライドサイズ
SLIDE_WIDTH = Inches(10)
SLIDE_HEIGHT = Inches(7.5)

# マージン・座標
LEFT_MARGIN = Emu(252000)        # 左マージン
RIGHT_EDGE = Emu(8748000)       # 右端（LEFT_MARGIN + 8496000）
CONTENT_WIDTH = Emu(8496000)    # コンテンツ幅
TITLE_TOP = Emu(164448)         # タイトルY座標
TITLE_HEIGHT = Emu(450932)      # タイトル高さ
HEADLINE_TOP = Emu(691977)      # ヘッドラインY座標
SECTION_BAR_TOP = Emu(1214454)  # セクションバーY座標
SECTION_BAR_HEIGHT = Emu(241556)
CONTENT_TOP = Emu(1506488)      # コンテンツ開始Y座標
SOURCE_TOP = Emu(6330701)       # Source行Y座標
SOURCE_HEIGHT = Emu(338075)
PAGE_NUM_TOP = Emu(6588540)     # ページ番号Y座標
PAGE_NUM_LEFT = Emu(5978455)
PAGE_NUM_WIDTH = Emu(3123808)
PAGE_NUM_HEIGHT = Emu(248302)

# カラー
BLUE = RGBColor(0x44, 0x72, 0xC4)      # セクションバー・ヘッダー
DARK_BLUE = RGBColor(0x20, 0x3A, 0x6B)  # テーブルヘッダー
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
BLACK = RGBColor(0x00, 0x00, 0x00)
DARK_GRAY = RGBColor(0x33, 0x33, 0x33)
LIGHT_GRAY = RGBColor(0xD6, 0xE4, 0xF0)  # テーブル交互行
VERY_LIGHT_GRAY = RGBColor(0xF2, 0xF2, 0xF2)

# フォント
FONT_EN = "Arial"
FONT_JA = "Meiryo"


# ---------------------------------------------------------------------------
# Utility Functions
# ---------------------------------------------------------------------------

def _set_font(run, size=Pt(9), bold=False, color=DARK_GRAY, font_name=FONT_EN):
    """runのフォントプロパティを設定。"""
    run.font.size = size
    run.font.bold = bold
    run.font.color.rgb = color
    run.font.name = font_name


def _add_textbox(slide, left, top, width, height, text="",
                 font_size=Pt(9), bold=False, color=DARK_GRAY,
                 font_name=FONT_EN, alignment=PP_ALIGN.LEFT):
    """テキストボックスを追加。"""
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = alignment
    run = p.add_run()
    run.text = text
    _set_font(run, size=font_size, bold=bold, color=color, font_name=font_name)
    return txBox


def _add_section_bar(slide, left, top, width, height, text="", color=BLUE):
    """セクションバー（青い帯+白文字）を追加。"""
    shape = slide.shapes.add_shape(
        1,  # MSO_SHAPE.RECTANGLE
        left, top, width, height
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()

    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = Emu(72000)
    tf.margin_top = Emu(0)
    tf.margin_bottom = Emu(0)
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    p.space_before = Pt(0)
    p.space_after = Pt(0)
    run = p.add_run()
    run.text = text
    _set_font(run, size=Pt(9), bold=True, color=WHITE, font_name=FONT_EN)
    return shape


def _add_title(slide, text, page_num=None):
    """スライドタイトルとページ番号を追加。"""
    _add_textbox(
        slide, LEFT_MARGIN, TITLE_TOP, CONTENT_WIDTH, TITLE_HEIGHT,
        text=text, font_size=Pt(16), bold=True, color=DARK_GRAY, font_name=FONT_EN,
    )
    if page_num is not None:
        _add_textbox(
            slide, PAGE_NUM_LEFT, PAGE_NUM_TOP, PAGE_NUM_WIDTH, PAGE_NUM_HEIGHT,
            text=f"| {page_num}", font_size=Pt(8), color=DARK_GRAY,
            alignment=PP_ALIGN.RIGHT,
        )


def _add_source(slide, text="Source:\tPublic filings, Kaisha Shikiho"):
    """Source行を追加。"""
    _add_textbox(
        slide, LEFT_MARGIN, SOURCE_TOP, CONTENT_WIDTH, SOURCE_HEIGHT,
        text=text, font_size=Pt(7), color=DARK_GRAY,
    )


def _set_cell(cell, text, font_size=Pt(8), bold=False, color=DARK_GRAY,
              font_name=FONT_EN, alignment=PP_ALIGN.LEFT, bg_color=None):
    """テーブルセルのテキストとスタイルを設定。"""
    cell.text = ""
    p = cell.text_frame.paragraphs[0]
    p.alignment = alignment
    p.space_before = Pt(0)
    p.space_after = Pt(0)
    run = p.add_run()
    run.text = str(text)
    _set_font(run, size=font_size, bold=bold, color=color, font_name=font_name)

    cell.margin_left = Emu(36000)
    cell.margin_right = Emu(36000)
    cell.margin_top = Emu(18000)
    cell.margin_bottom = Emu(18000)
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE

    if bg_color:
        cell.fill.solid()
        cell.fill.fore_color.rgb = bg_color


def _fmt_num(val, fmt="comma"):
    """数値フォーマット。"""
    if val is None:
        return "N/A"
    if fmt == "comma":
        return f"{val:,.0f}"
    elif fmt == "pct":
        return f"{val:.1%}"
    elif fmt == "x":
        return f"{val:.1f}x"
    elif fmt == "pct_simple":
        return f"{val * 100:.1f}%"
    return str(val)


# ---------------------------------------------------------------------------
# Slide 1: Company Overview
# ---------------------------------------------------------------------------

def _build_overview_slide(prs, profile, page_num=1):
    """Company Overview スライドを生成。"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout

    company_name_en = profile.get('company_name_en') or profile.get('company_name', '')
    code = profile.get('code', '')

    # タイトル
    _add_title(slide, f"{company_name_en} : Company Overview", page_num)

    # ヘッドライン
    web = profile.get('web', {})
    headline = web.get('headline_en', '')
    if headline:
        _add_textbox(
            slide, LEFT_MARGIN, HEADLINE_TOP, CONTENT_WIDTH, Emu(400000),
            text=headline, font_size=Pt(9), color=DARK_GRAY,
        )

    # --- セクションバー: Company Info ---
    _add_section_bar(
        slide, LEFT_MARGIN, SECTION_BAR_TOP,
        Emu(2700000), SECTION_BAR_HEIGHT,
        text="Company Info"
    )

    # --- Company Info テーブル ---
    fin = profile.get('financial', {})
    rows_data = [
        ("Founded", web.get('founding_year', 'N/A')),
        ("Headquarters", web.get('headquarters', 'N/A')),
        ("Employees", str(profile.get('num_employees', 'N/A'))),
        ("Representative", profile.get('representative', 'N/A')),
        ("Main Business", web.get('main_business_en', 'N/A')[:80]),
        ("Global Footprint", web.get('global_footprint', 'N/A')[:80]),
        ("Group Companies", web.get('group_companies', 'N/A')[:60]),
    ]

    table_height = Emu(int(2200000 * len(rows_data) / 7))
    table = slide.shapes.add_table(
        len(rows_data), 2,
        LEFT_MARGIN, CONTENT_TOP,
        Emu(2700000), table_height
    ).table

    # 列幅
    table.columns[0].width = Emu(900000)
    table.columns[1].width = Emu(1800000)

    for i, (label, value) in enumerate(rows_data):
        bg = VERY_LIGHT_GRAY if i % 2 == 0 else None
        _set_cell(table.cell(i, 0), label, bold=True, font_size=Pt(7),
                  bg_color=bg)
        _set_cell(table.cell(i, 1), value, font_size=Pt(7), bg_color=bg)

    # --- セクションバー: Shareholder Composition ---
    sh_bar_top = Emu(3756419)
    _add_section_bar(
        slide, LEFT_MARGIN, sh_bar_top,
        Emu(2700000), SECTION_BAR_HEIGHT,
        text="Shareholder Composition"
    )

    # --- 株主テーブル ---
    shareholders = profile.get('shareholders', [])[:10]
    if shareholders:
        sh_top = Emu(4048453)
        sh_rows = len(shareholders) + 1  # +header
        sh_table = slide.shapes.add_table(
            sh_rows, 3,
            LEFT_MARGIN, sh_top,
            Emu(2700000), Emu(min(2100000, sh_rows * 180000))
        ).table

        sh_table.columns[0].width = Emu(250000)
        sh_table.columns[1].width = Emu(1850000)
        sh_table.columns[2].width = Emu(600000)

        # Header
        for j, hdr in enumerate(["#", "Shareholder", "Ratio"]):
            _set_cell(sh_table.cell(0, j), hdr, bold=True, font_size=Pt(6),
                      color=WHITE, bg_color=DARK_BLUE,
                      alignment=PP_ALIGN.CENTER if j != 1 else PP_ALIGN.LEFT)

        for i, sh in enumerate(shareholders):
            row_idx = i + 1
            bg = VERY_LIGHT_GRAY if i % 2 == 0 else None
            _set_cell(sh_table.cell(row_idx, 0), str(sh.get('rank', i + 1)),
                      font_size=Pt(6), alignment=PP_ALIGN.CENTER, bg_color=bg)
            _set_cell(sh_table.cell(row_idx, 1), sh.get('name', '')[:40],
                      font_size=Pt(6), font_name=FONT_JA, bg_color=bg)
            ratio = sh.get('ratio')
            ratio_str = f"{ratio:.1%}" if ratio else "N/A"
            _set_cell(sh_table.cell(row_idx, 2), ratio_str,
                      font_size=Pt(6), alignment=PP_ALIGN.RIGHT, bg_color=bg)

    # --- 株価チャート（右上エリア） ---
    stock_history = profile.get('stock_history', [])
    if stock_history:
        chart_path = _generate_stock_chart(stock_history, code, company_name_en)
        if chart_path:
            slide.shapes.add_picture(
                chart_path, Emu(6048000), CONTENT_TOP,
                Emu(2700000), Emu(2200000)
            )
            # セクションバー
            _add_section_bar(
                slide, Emu(6048000), sh_bar_top,
                Emu(2700000), SECTION_BAR_HEIGHT,
                text=f"Share Price (TSE: {code})"
            )

    # --- 売上・利益チャート（右中央エリア） ---
    _add_section_bar(
        slide, Emu(3150000), SECTION_BAR_TOP,
        Emu(2700000), SECTION_BAR_HEIGHT,
        text="Revenue & Operating Income"
    )

    # Revenue/OI chart
    rev_chart_path = _generate_revenue_chart(profile)
    if rev_chart_path:
        slide.shapes.add_picture(
            rev_chart_path, Emu(3150000), CONTENT_TOP,
            Emu(2700000), Emu(2200000)
        )

    # Source
    _add_source(slide)

    return slide


def _generate_stock_chart(stock_history, code, company_name=""):
    """matplotlib で株価チャートを生成 → PNG パスを返す。"""
    try:
        import matplotlib
        matplotlib.use('Agg')
        import matplotlib.pyplot as plt
        import matplotlib.dates as mdates
        from datetime import datetime

        dates = [datetime.strptime(d, "%Y-%m-%d") for d, _ in stock_history]
        prices = [p for _, p in stock_history]

        fig, ax = plt.subplots(figsize=(5, 3.5))
        ax.plot(dates, prices, color='#4472C4', linewidth=1.2)
        ax.fill_between(dates, prices, alpha=0.1, color='#4472C4')
        ax.set_ylabel('JPY', fontsize=8)
        ax.xaxis.set_major_formatter(mdates.DateFormatter('%b %y'))
        ax.xaxis.set_major_locator(mdates.MonthLocator(interval=3))
        ax.tick_params(labelsize=7)
        ax.grid(True, alpha=0.3)
        ax.spines['top'].set_visible(False)
        ax.spines['right'].set_visible(False)

        # 最新株価を表示
        latest_price = prices[-1]
        ax.annotate(
            f'¥{latest_price:,.0f}',
            xy=(dates[-1], latest_price),
            fontsize=8, fontweight='bold', color='#4472C4',
            ha='right', va='bottom',
        )

        plt.tight_layout()
        tmp = tempfile.NamedTemporaryFile(suffix='.png', delete=False)
        fig.savefig(tmp.name, dpi=150, bbox_inches='tight', transparent=True)
        plt.close(fig)
        return tmp.name
    except Exception as e:
        print(f"  [Chart] 株価チャート生成失敗: {e}")
        return None


def _generate_revenue_chart(profile):
    """売上高・営業利益の棒グラフを生成。"""
    try:
        import matplotlib
        matplotlib.use('Agg')
        import matplotlib.pyplot as plt

        fin = profile.get('financial', {})
        rev_ltm = fin.get('rev_ltm')
        op_ltm = fin.get('op_ltm')
        rev_fwd = fin.get('rev_forecast')
        op_fwd = fin.get('op_forecast')

        if not rev_ltm:
            return None

        labels = []
        revs = []
        ops = []

        labels.append('LTM')
        revs.append(rev_ltm)
        ops.append(op_ltm or 0)

        if rev_fwd:
            labels.append('FY E')
            revs.append(rev_fwd)
            ops.append(op_fwd or 0)

        fig, ax = plt.subplots(figsize=(5, 3.5))
        x = range(len(labels))
        width = 0.35

        bars1 = ax.bar([i - width / 2 for i in x], revs, width,
                       label='Revenue', color='#4472C4', alpha=0.8)
        bars2 = ax.bar([i + width / 2 for i in x], ops, width,
                       label='Op. Income', color='#ED7D31', alpha=0.8)

        ax.set_xticks(list(x))
        ax.set_xticklabels(labels, fontsize=8)
        ax.set_ylabel('JPY mn', fontsize=8)
        ax.tick_params(labelsize=7)
        ax.legend(fontsize=7, loc='upper right')
        ax.grid(True, axis='y', alpha=0.3)
        ax.spines['top'].set_visible(False)
        ax.spines['right'].set_visible(False)

        # 値ラベル
        for bar in bars1:
            ax.text(bar.get_x() + bar.get_width() / 2, bar.get_height(),
                    f'{bar.get_height():,.0f}', ha='center', va='bottom', fontsize=6)
        for bar in bars2:
            ax.text(bar.get_x() + bar.get_width() / 2, bar.get_height(),
                    f'{bar.get_height():,.0f}', ha='center', va='bottom', fontsize=6)

        plt.tight_layout()
        tmp = tempfile.NamedTemporaryFile(suffix='.png', delete=False)
        fig.savefig(tmp.name, dpi=150, bbox_inches='tight', transparent=True)
        plt.close(fig)
        return tmp.name
    except Exception as e:
        print(f"  [Chart] 売上チャート生成失敗: {e}")
        return None


# ---------------------------------------------------------------------------
# Slides 5-6: Board of Directors
# ---------------------------------------------------------------------------

def _build_directors_slides(prs, profile, start_page=2):
    """役員スライド（6人/スライド、動的ページネーション）を生成。"""
    directors = profile.get('directors', [])
    if not directors:
        return start_page

    company_name_en = profile.get('company_name_en') or profile.get('company_name', '')
    directors_per_page = 6  # 3 left + 3 right
    pages = (len(directors) + directors_per_page - 1) // directors_per_page
    page_num = start_page

    for page_idx in range(pages):
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        page_label = f"({page_idx + 1}/{pages})" if pages > 1 else ""
        _add_title(
            slide,
            f"{company_name_en} : Board of Directors {page_label}",
            page_num,
        )

        # Independent member badge
        badge = slide.shapes.add_shape(
            5,  # ROUNDED_RECTANGLE
            Emu(6370480), Emu(296718), Emu(2376264), Emu(201108)
        )
        badge.fill.solid()
        badge.fill.fore_color.rgb = BLUE
        badge.line.fill.background()
        tf = badge.text_frame
        tf.margin_left = Emu(36000)
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        run = p.add_run()
        run.text = "Independent/external member"
        _set_font(run, size=Pt(7), bold=False, color=WHITE)

        # 左テーブル（3人）と右テーブル（3人）
        start_idx = page_idx * directors_per_page
        left_dirs = directors[start_idx:start_idx + 3]
        right_dirs = directors[start_idx + 3:start_idx + 6]

        for col_idx, dirs in enumerate([left_dirs, right_dirs]):
            if not dirs:
                continue

            table_left = LEFT_MARGIN if col_idx == 0 else Emu(4570744)
            table_top = Emu(789087)
            table_width = Emu(4176000)

            # テーブル生成: rows = ヘッダー + 各役員2行(名前+経歴)
            num_rows = 1 + len(dirs) * 2  # header + (name_row + career_row) * N
            table = slide.shapes.add_table(
                num_rows, 2,
                table_left, table_top,
                table_width, Emu(5400000)
            ).table

            table.columns[0].width = Emu(1200000)
            table.columns[1].width = Emu(2976000)

            # Header
            _set_cell(table.cell(0, 0), "Name / Title", bold=True,
                      font_size=Pt(7), color=WHITE, bg_color=DARK_BLUE)
            _set_cell(table.cell(0, 1), "Career Summary", bold=True,
                      font_size=Pt(7), color=WHITE, bg_color=DARK_BLUE)

            for i, d in enumerate(dirs):
                name_row = 1 + i * 2
                career_row = name_row + 1

                # 名前・役職
                name_text = d.get('name', '')
                title_text = d.get('title', '')
                dob = d.get('dob', '')
                name_display = f"{name_text}\n{title_text}"
                if dob:
                    name_display += f"\n({dob})"

                is_independent = any(kw in title_text for kw in
                                     ['社外', '独立', 'Independent', 'External'])
                name_color = BLUE if is_independent else DARK_GRAY

                _set_cell(table.cell(name_row, 0), name_display,
                          font_size=Pt(7), font_name=FONT_JA, color=name_color,
                          bold=True)

                # 経歴
                career = d.get('career', '')
                # HTMLタグ除去
                career = re.sub(r'<[^>]+>', '', career)
                career = career[:500]  # 長すぎる経歴を切り詰め

                _set_cell(table.cell(name_row, 1), career,
                          font_size=Pt(6), font_name=FONT_JA)

                # 空行（スペーサー）
                _set_cell(table.cell(career_row, 0), "")
                _set_cell(table.cell(career_row, 1), "")
                table.rows[career_row].height = Emu(36000)

        _add_source(slide, "Source:\tPublic filings (EDINET)")
        page_num += 1

    return page_num


# ---------------------------------------------------------------------------
# Slides 7-8: Comparable Companies Analysis
# ---------------------------------------------------------------------------

def _build_comps_slides(prs, comps_data, target_code, company_name_en, start_page=4):
    """Comps分析スライド（2ページ: BS/P&L + Valuation）を生成。"""
    if not comps_data:
        return start_page

    page_num = start_page

    # --- Slide 7: Company Info + BS + P&L ---
    slide1 = prs.slides.add_slide(prs.slide_layouts[6])
    _add_title(
        slide1,
        f"{company_name_en} : Comparable Companies Analysis (1/2)",
        page_num,
    )

    # ヘッドライン
    _add_textbox(
        slide1, LEFT_MARGIN, HEADLINE_TOP, CONTENT_WIDTH, Emu(270000),
        text="Selected public comparable companies trading multiples",
        font_size=Pt(9), color=DARK_GRAY,
    )

    # テーブル構築
    headers1 = [
        "Company", "Code", "Mkt Cap\n(JPY mn)", "EV\n(JPY mn)",
        "Rev LTM\n(JPY mn)", "OP LTM\n(JPY mn)", "EBITDA LTM\n(JPY mn)",
        "Cash\n(JPY mn)", "Debt\n(JPY mn)",
    ]

    num_companies = len(comps_data)
    # +1 header, +1 spacer for target, +4 for stats
    num_rows = 1 + num_companies + 4

    table1 = slide1.shapes.add_table(
        num_rows, len(headers1),
        LEFT_MARGIN, Emu(1100000),
        CONTENT_WIDTH, Emu(4800000)
    ).table

    # 列幅設定
    col_widths = [1800000, 600000, 900000, 900000, 900000, 900000, 900000, 700000, 700000]
    for j, w in enumerate(col_widths):
        table1.columns[j].width = Emu(w)

    # Header
    for j, hdr in enumerate(headers1):
        _set_cell(table1.cell(0, j), hdr, bold=True, font_size=Pt(7),
                  color=WHITE, bg_color=DARK_BLUE,
                  alignment=PP_ALIGN.CENTER if j > 0 else PP_ALIGN.LEFT)

    # Data rows
    values_for_stats = {h: [] for h in headers1[2:]}

    for i, cd in enumerate(comps_data):
        row = i + 1
        is_target = str(cd.get('code', '')) == str(target_code)
        bg = LIGHT_GRAY if is_target else (VERY_LIGHT_GRAY if i % 2 == 0 else None)
        name_bold = is_target

        _set_cell(table1.cell(row, 0), cd.get('name', '')[:25],
                  font_size=Pt(7), font_name=FONT_JA, bold=name_bold, bg_color=bg)
        _set_cell(table1.cell(row, 1), str(cd.get('code', '')),
                  font_size=Pt(7), alignment=PP_ALIGN.CENTER, bg_color=bg)

        num_fields = [
            ('market_cap', 'comma'), ('_ev', 'comma'),
            ('rev_ltm', 'comma'), ('op_ltm', 'comma'),
            ('ebitda_ltm', 'comma'), ('cash', 'comma'),
            ('total_debt', 'comma'),
        ]

        for j, (key, fmt) in enumerate(num_fields):
            val = cd.get(key)
            if val is None and key == '_ev':
                from financial_calc import calc_ev
                val = calc_ev(cd.get('market_cap'), cd.get('total_debt'), cd.get('cash'))
            _set_cell(table1.cell(row, j + 2), _fmt_num(val, fmt),
                      font_size=Pt(7), alignment=PP_ALIGN.RIGHT, bg_color=bg)
            if val is not None and not is_target:
                values_for_stats[headers1[j + 2]].append(val)

    # Stats rows
    stat_start = num_companies + 1
    for s_idx, (stat_name, stat_fn) in enumerate([
        ("Median", lambda vs: sorted(vs)[len(vs) // 2] if vs else None),
        ("Mean", lambda vs: sum(vs) / len(vs) if vs else None),
        ("Low", lambda vs: min(vs) if vs else None),
        ("High", lambda vs: max(vs) if vs else None),
    ]):
        row = stat_start + s_idx
        _set_cell(table1.cell(row, 0), stat_name, bold=True, font_size=Pt(7),
                  bg_color=LIGHT_GRAY)
        _set_cell(table1.cell(row, 1), "", bg_color=LIGHT_GRAY)
        for j, hdr in enumerate(headers1[2:]):
            vals = values_for_stats.get(hdr, [])
            val = stat_fn(vals)
            _set_cell(table1.cell(row, j + 2), _fmt_num(val, 'comma'),
                      font_size=Pt(7), alignment=PP_ALIGN.RIGHT, bg_color=LIGHT_GRAY)

    _add_source(slide1)
    page_num += 1

    # --- Slide 8: Valuation Multiples ---
    slide2 = prs.slides.add_slide(prs.slide_layouts[6])
    _add_title(
        slide2,
        f"{company_name_en} : Comparable Companies Analysis (2/2)",
        page_num,
    )

    headers2 = [
        "Company", "Code", "EV/EBITDA\nLTM", "EV/EBITDA\nFwd",
        "PER\nFwd", "PBR", "Div Yield",
        "OP Margin\nLTM", "EBITDA Margin\nLTM",
    ]

    num_rows2 = 1 + num_companies + 4
    table2 = slide2.shapes.add_table(
        num_rows2, len(headers2),
        LEFT_MARGIN, Emu(1100000),
        CONTENT_WIDTH, Emu(4800000)
    ).table

    col_widths2 = [1800000, 600000, 900000, 900000, 900000, 800000, 800000, 900000, 900000]
    for j, w in enumerate(col_widths2):
        table2.columns[j].width = Emu(w)

    for j, hdr in enumerate(headers2):
        _set_cell(table2.cell(0, j), hdr, bold=True, font_size=Pt(7),
                  color=WHITE, bg_color=DARK_BLUE,
                  alignment=PP_ALIGN.CENTER if j > 0 else PP_ALIGN.LEFT)

    mult_stats = {h: [] for h in headers2[2:]}

    for i, cd in enumerate(comps_data):
        row = i + 1
        is_target = str(cd.get('code', '')) == str(target_code)
        bg = LIGHT_GRAY if is_target else (VERY_LIGHT_GRAY if i % 2 == 0 else None)

        _set_cell(table2.cell(row, 0), cd.get('name', '')[:25],
                  font_size=Pt(7), font_name=FONT_JA, bold=is_target, bg_color=bg)
        _set_cell(table2.cell(row, 1), str(cd.get('code', '')),
                  font_size=Pt(7), alignment=PP_ALIGN.CENTER, bg_color=bg)

        multiples = cd.get('_multiples', {})
        from financial_calc import safe_div, calc_ev

        ev = cd.get('_ev')
        if ev is None:
            ev = calc_ev(cd.get('market_cap'), cd.get('total_debt'), cd.get('cash'))

        mult_vals = [
            (multiples.get('ev_ebitda_ltm') or safe_div(ev, cd.get('ebitda_ltm')), 'x'),
            (multiples.get('ev_ebitda_fwd') or safe_div(ev, cd.get('ebitda_forecast')), 'x'),
            (multiples.get('per_fwd') or safe_div(cd.get('market_cap'), cd.get('ni_forecast')), 'x'),
            (multiples.get('pbr') or safe_div(cd.get('market_cap'), cd.get('equity_parent')), 'x'),
            (multiples.get('div_yield') or safe_div(cd.get('dps'), cd.get('stock_price')), 'pct_simple'),
            (safe_div(cd.get('op_ltm'), cd.get('rev_ltm')), 'pct_simple'),
            (safe_div(cd.get('ebitda_ltm'), cd.get('rev_ltm')), 'pct_simple'),
        ]

        for j, (val, fmt) in enumerate(mult_vals):
            _set_cell(table2.cell(row, j + 2), _fmt_num(val, fmt),
                      font_size=Pt(7), alignment=PP_ALIGN.RIGHT, bg_color=bg)
            if val is not None and not is_target:
                mult_stats[headers2[j + 2]].append(val)

    # Stats
    stat_start2 = num_companies + 1
    for s_idx, (stat_name, stat_fn) in enumerate([
        ("Median", lambda vs: sorted(vs)[len(vs) // 2] if vs else None),
        ("Mean", lambda vs: sum(vs) / len(vs) if vs else None),
        ("Low", lambda vs: min(vs) if vs else None),
        ("High", lambda vs: max(vs) if vs else None),
    ]):
        row = stat_start2 + s_idx
        _set_cell(table2.cell(row, 0), stat_name, bold=True, font_size=Pt(7),
                  bg_color=LIGHT_GRAY)
        _set_cell(table2.cell(row, 1), "", bg_color=LIGHT_GRAY)
        for j, hdr in enumerate(headers2[2:]):
            vals = mult_stats.get(hdr, [])
            val = stat_fn(vals)
            fmt = 'pct_simple' if 'Margin' in hdr or 'Yield' in hdr else 'x'
            _set_cell(table2.cell(row, j + 2), _fmt_num(val, fmt),
                      font_size=Pt(7), alignment=PP_ALIGN.RIGHT, bg_color=LIGHT_GRAY)

    _add_source(slide2)
    page_num += 1

    return page_num


# ---------------------------------------------------------------------------
# Slide 9: Preliminary Financial Analysis
# ---------------------------------------------------------------------------

def _build_financial_slide(prs, profile, comps_data, start_page=6):
    """Preliminary Financial Analysis スライドを生成。"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    company_name_en = profile.get('company_name_en') or profile.get('company_name', '')

    _add_title(
        slide,
        f"{company_name_en} : Preliminary Financial Analysis",
        start_page,
    )

    fin = profile.get('financial', {})

    # Detailed P&L table
    headers = ["", "LTM", "FY Forecast"]
    rows_data = [
        ("Revenue (JPY mn)", fin.get('rev_ltm'), fin.get('rev_forecast')),
        ("Operating Income (JPY mn)", fin.get('op_ltm'), fin.get('op_forecast')),
        ("Net Income (JPY mn)", fin.get('ni_ltm'), fin.get('ni_forecast')),
        ("D&A (JPY mn)", fin.get('da_ltm'), fin.get('da_ltm')),
        ("EBITDA (JPY mn)", fin.get('ebitda_ltm'), fin.get('ebitda_forecast')),
        ("", "", ""),
        ("OP Margin", _safe_margin(fin.get('op_ltm'), fin.get('rev_ltm')),
         _safe_margin(fin.get('op_forecast'), fin.get('rev_forecast'))),
        ("EBITDA Margin", _safe_margin(fin.get('ebitda_ltm'), fin.get('rev_ltm')),
         _safe_margin(fin.get('ebitda_forecast'), fin.get('rev_forecast'))),
        ("", "", ""),
        ("Stock Price (JPY)", fin.get('stock_price'), ""),
        ("Shares Outstanding (K)", fin.get('shares_outstanding'), ""),
        ("Market Cap (JPY mn)", fin.get('market_cap'), ""),
        ("Cash (JPY mn)", fin.get('cash'), ""),
        ("Total Debt (JPY mn)", fin.get('total_debt'), ""),
        ("EV (JPY mn)", fin.get('_ev'), ""),
        ("", "", ""),
        ("EV/EBITDA LTM", "", ""),
        ("EV/EBITDA Fwd", "", ""),
        ("PER Fwd", "", ""),
        ("PBR", "", ""),
        ("Dividend Yield", "", ""),
    ]

    # Fill in multiples
    mult = fin.get('_multiples', {})
    rows_data[16] = ("EV/EBITDA LTM", _fmt_mult(mult.get('ev_ebitda_ltm'), 'x'), "")
    rows_data[17] = ("EV/EBITDA Fwd", _fmt_mult(mult.get('ev_ebitda_fwd'), 'x'), "")
    rows_data[18] = ("PER Fwd", _fmt_mult(mult.get('per_fwd'), 'x'), "")
    rows_data[19] = ("PBR", _fmt_mult(mult.get('pbr'), 'x'), "")
    rows_data[20] = ("Dividend Yield", _fmt_mult(mult.get('div_yield'), 'pct'), "")

    table = slide.shapes.add_table(
        len(rows_data), 3,
        LEFT_MARGIN, Emu(1000000),
        Emu(4500000), Emu(5200000)
    ).table

    table.columns[0].width = Emu(2200000)
    table.columns[1].width = Emu(1150000)
    table.columns[2].width = Emu(1150000)

    # Header
    for j, hdr in enumerate(headers):
        _set_cell(table.cell(0, j), hdr, bold=True, font_size=Pt(7),
                  color=WHITE, bg_color=DARK_BLUE,
                  alignment=PP_ALIGN.CENTER if j > 0 else PP_ALIGN.LEFT)

    for i, (label, val1, val2) in enumerate(rows_data):
        row = i  # rows_data includes header conceptually but we add separately
        if i == 0:
            continue  # skip, already handled
        actual_row = i

        is_separator = label == ""
        bg = None
        if is_separator:
            _set_cell(table.cell(actual_row, 0), "", font_size=Pt(4))
            _set_cell(table.cell(actual_row, 1), "", font_size=Pt(4))
            _set_cell(table.cell(actual_row, 2), "", font_size=Pt(4))
            table.rows[actual_row].height = Emu(72000)
            continue

        _set_cell(table.cell(actual_row, 0), label, font_size=Pt(7), bold=True)

        for j, val in enumerate([val1, val2], start=1):
            if isinstance(val, str):
                _set_cell(table.cell(actual_row, j), val,
                          font_size=Pt(7), alignment=PP_ALIGN.RIGHT)
            elif val is not None:
                _set_cell(table.cell(actual_row, j), _fmt_num(val),
                          font_size=Pt(7), alignment=PP_ALIGN.RIGHT)
            else:
                _set_cell(table.cell(actual_row, j), "",
                          font_size=Pt(7), alignment=PP_ALIGN.RIGHT)

    # Implied valuation range (if comps data available)
    if comps_data:
        comp_multiples = []
        for cd in comps_data:
            if str(cd.get('code', '')) == str(profile.get('code', '')):
                continue
            m = cd.get('_multiples', {})
            ev_ebitda = m.get('ev_ebitda_ltm')
            if ev_ebitda:
                comp_multiples.append(ev_ebitda)

        if comp_multiples and fin.get('ebitda_ltm') and fin.get('cash') is not None:
            median_mult = sorted(comp_multiples)[len(comp_multiples) // 2]
            low_mult = min(comp_multiples)
            high_mult = max(comp_multiples)

            ebitda = fin['ebitda_ltm']
            cash = fin.get('cash', 0) or 0
            debt = fin.get('total_debt', 0) or 0
            shares = fin.get('shares_outstanding')

            _add_textbox(
                slide, Emu(5000000), Emu(1000000), Emu(3700000), Emu(300000),
                text="Implied Valuation Range (EV/EBITDA LTM)",
                font_size=Pt(10), bold=True, color=DARK_BLUE,
            )

            val_headers = ["", "Low", "Median", "High"]
            val_rows = [
                ("EV/EBITDA Multiple",
                 f"{low_mult:.1f}x", f"{median_mult:.1f}x", f"{high_mult:.1f}x"),
                ("Implied EV (JPY mn)",
                 f"{ebitda * low_mult:,.0f}", f"{ebitda * median_mult:,.0f}",
                 f"{ebitda * high_mult:,.0f}"),
                ("- Net Debt (JPY mn)",
                 f"{debt - cash:,.0f}", f"{debt - cash:,.0f}", f"{debt - cash:,.0f}"),
                ("Implied Eq. Value",
                 f"{ebitda * low_mult - debt + cash:,.0f}",
                 f"{ebitda * median_mult - debt + cash:,.0f}",
                 f"{ebitda * high_mult - debt + cash:,.0f}"),
            ]

            if shares and shares > 0:
                val_rows.append((
                    "Implied Share Price",
                    f"¥{(ebitda * low_mult - debt + cash) / shares * 1000:,.0f}",
                    f"¥{(ebitda * median_mult - debt + cash) / shares * 1000:,.0f}",
                    f"¥{(ebitda * high_mult - debt + cash) / shares * 1000:,.0f}",
                ))

            val_table = slide.shapes.add_table(
                len(val_rows) + 1, 4,
                Emu(5000000), Emu(1400000),
                Emu(3700000), Emu(2000000)
            ).table

            val_table.columns[0].width = Emu(1500000)
            val_table.columns[1].width = Emu(750000)
            val_table.columns[2].width = Emu(750000)
            val_table.columns[3].width = Emu(700000)

            for j, hdr in enumerate(val_headers):
                _set_cell(val_table.cell(0, j), hdr, bold=True, font_size=Pt(7),
                          color=WHITE, bg_color=DARK_BLUE,
                          alignment=PP_ALIGN.CENTER if j > 0 else PP_ALIGN.LEFT)

            for i, (label, *vals) in enumerate(val_rows):
                row = i + 1
                is_last = i == len(val_rows) - 1
                _set_cell(val_table.cell(row, 0), label,
                          font_size=Pt(7), bold=is_last)
                for j, v in enumerate(vals):
                    _set_cell(val_table.cell(row, j + 1), v,
                              font_size=Pt(7), alignment=PP_ALIGN.RIGHT,
                              bold=is_last,
                              bg_color=LIGHT_GRAY if is_last else None)

    _add_source(slide)
    return start_page + 1


def _safe_margin(num, denom):
    """マージン計算。"""
    if num is None or denom is None or denom == 0:
        return None
    return num / denom


def _fmt_mult(val, fmt):
    """マルチプルフォーマット。"""
    if val is None:
        return "N/A"
    if fmt == 'x':
        return f"{val:.1f}x"
    if fmt == 'pct':
        return f"{val:.1%}"
    return str(val)


# ---------------------------------------------------------------------------
# Main Builder
# ---------------------------------------------------------------------------

def build_profile_pptx(profile, comps_data=None, output_path=None):
    """
    Company Profile PPTX を生成。

    Parameters:
        profile: collect_profile_data() の戻り値
        comps_data: list of build_company_data() 戻り値（Comps対象企業）
        output_path: 出力先パス（Noneの場合デフォルト名）

    Returns:
        str: 出力ファイルパス
    """
    code = profile.get('code', 'unknown')
    company_name_en = profile.get('company_name_en') or profile.get('company_name', '')

    if output_path is None:
        safe_name = re.sub(r'[^\w\s-]', '', company_name_en).strip().replace(' ', '_')
        output_path = f"{code}_{safe_name}_Company_Profile.pptx"

    print(f"\n=== PPTX生成開始: {output_path} ===")

    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT

    page = 1

    # Slide 1: Company Overview
    _build_overview_slide(prs, profile, page)
    page += 1

    # Slides 2-N: Board of Directors
    page = _build_directors_slides(prs, profile, page)

    # Slides N+1, N+2: Comps Analysis
    if comps_data:
        page = _build_comps_slides(
            prs, comps_data, code, company_name_en, page
        )

    # Final Slide: Financial Analysis
    _build_financial_slide(prs, profile, comps_data, page)

    prs.save(output_path)
    print(f"=== PPTX保存完了: {output_path} ===")

    # 一時チャート画像をクリーンアップ
    import glob
    for tmp in glob.glob(os.path.join(tempfile.gettempdir(), "tmp*.png")):
        try:
            os.remove(tmp)
        except Exception:
            pass

    return output_path


if __name__ == "__main__":
    # テスト用: ダミーデータで生成
    profile = {
        'code': '6763',
        'company_name': '帝国通信工業',
        'company_name_en': 'Teikoku Tsushin Kogyo',
        'representative': '代表取締役社長　羽生 益雄',
        'num_employees': 1586,
        'shareholders': [
            {'rank': 1, 'name': '日本マスタートラスト信託銀行', 'ratio': 0.1135},
            {'rank': 2, 'name': 'HSBC HONG KONG', 'ratio': 0.0762},
        ],
        'directors': [
            {'name': '羽生 益雄', 'title': '代表取締役社長', 'career': '1977年入社', 'dob': '1956-04-22'},
        ],
        'web': {
            'founding_year': '1944',
            'headquarters': 'Meguro-ku, Tokyo',
            'headline_en': 'Leading manufacturer of variable resistors and electronic components.',
            'main_business_en': 'Electronic components manufacturer.',
            'global_footprint': 'Japan, China',
            'group_companies': '16 consolidated subsidiaries',
        },
        'financial': {
            'stock_price': 2745,
            'shares_outstanding': 9374,
            'market_cap': 25710,
            'cash': 12006,
            'total_debt': 77,
            'equity_parent': 27494,
            'rev_ltm': 17316,
            'op_ltm': 1859,
            'ni_ltm': 1757,
            'da_ltm': 771,
            'ebitda_ltm': 2630,
            'rev_forecast': 16800,
            'op_forecast': 1300,
            'ni_forecast': 1200,
            'ebitda_forecast': 2058,
            'dps': 100,
            '_ev': 13881,
            '_multiples': {
                'ev_ebitda_ltm': 5.3,
                'ev_ebitda_fwd': 6.7,
                'per_fwd': 21.4,
                'pbr': 0.94,
                'div_yield': 0.036,
            },
        },
        'stock_history': [],
    }

    build_profile_pptx(profile, output_path="test_profile.pptx")
